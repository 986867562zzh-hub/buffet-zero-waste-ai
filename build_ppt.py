"""
Build the complete 17-slide presentation for Group Project
Smart Tourism Studies - AI Zero-Waste Buffet Management System
"""
import os
os.chdir(r'D:\Quant_Projects\Quant_Projects\buffet-demo')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import time

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ===== Color Scheme =====
DARK = RGBColor(0x2D, 0x34, 0x36)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT2 = RGBColor(0xFF, 0xA5, 0x02)
GREEN = RGBColor(0x27, 0xAE, 0x60)
BLUE = RGBColor(0x34, 0x98, 0xDB)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xFA)
MED_GRAY = RGBColor(0x95, 0xA5, 0xA6)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)

def add_bg(slide, color=LIGHT_GRAY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at top"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.font.bold = True

    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11), Inches(0.35))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = MED_GRAY

def add_body_text(slide, text, left=0.8, top=1.6, width=11.5, height=5, size=18, bold=False, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.line_spacing = Pt(size * 1.5)
    return tf

def add_card(slide, left, top, width, height, title, body, title_color=ACCENT):
    """Add a card-style box"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = title_color
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK
    p2.line_spacing = Pt(22)
    return shape

def add_page_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.RIGHT

# ================================================================
# SLIDE 1: Title Slide
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Big icon
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.0), Inches(2.3), Inches(2.3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = "🍽️"
p.font.size = Pt(80)
p.alignment = PP_ALIGN.CENTER

# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.3), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "AI-Powered Zero-Waste Buffet Management System"
p.font.size = Pt(40)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Subtitle
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.3), Inches(0.6))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "Smart Tourism Operations & Management — Group Project"
p2.font.size = Pt(20)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER

# Team members
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.3), Inches(1))
tf3 = txBox3.text_frame
members = ["张正浩", "Member 2", "Member 3", "Member 4", "Member 5"]
for i, m in enumerate(members):
    if i == 0:
        p3 = tf3.paragraphs[0]
    else:
        p3 = tf3.add_paragraph()
    p3.text = f"  {m}  " if i < len(members)-1 else f"  {m}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = WHITE
    p3.alignment = PP_ALIGN.CENTER

add_page_number(slide, 1)


# ================================================================
# SLIDE 2: Problem Statement
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "The Problem: Buffet Food Waste", "A $1 trillion global problem in the hospitality industry")

# Key stats cards
add_card(slide, 0.8, 1.6, 3.8, 2.5,
    "🌍 Global Scale",
    "• 1.3 billion tons of food wasted annually\n"
    "• Buffet restaurants waste 2× more than à la carte\n"
    "• 30-40% of buffet food goes uneaten\n"
    "• Hotel buffets generate ~200kg waste/day")

add_card(slide, 5.0, 1.6, 3.8, 2.5,
    "👤 Consumer Behavior",
    "• 'Eyes bigger than stomach' syndrome\n"
    "• Taking too much to 'get value for money'\n"
    "• No real-time feedback on over-consumption\n"
    "• Plate waste ends up in landfills")

add_card(slide, 9.2, 1.6, 3.8, 2.5,
    "🏨 Hotel Pain Point",
    "• Over-preparation: cooking for worst case\n"
    "• Leftover buffet dishes hard to re-sell\n"
    "• High food cost + waste disposal cost\n"
    "• Sustainability pressure from guests & regulators")

# Bottom: the core question
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "How can AI help reduce food waste while improving the dining experience?"
p.font.size = Pt(24)
p.font.color.rgb = ACCENT
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "A small, specific problem — with a smart, AI-powered solution."
p2.font.size = Pt(16)
p2.font.color.rgb = MED_GRAY
p2.alignment = PP_ALIGN.CENTER

add_page_number(slide, 2)


# ================================================================
# SLIDE 3: User & Scenario
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Target Users & Scenario", "Integrated resort buffet restaurant — Macau context")

# Two user personas
add_card(slide, 0.8, 1.6, 5.8, 3.5,
    "🍽️ Primary User: Diner (Consumer)",
    "• Tourist or business traveler at hotel buffet\n"
    "• Wants to enjoy variety but tends to over-take\n"
    "• Unaware of cumulative calorie intake\n"
    "• Motivated by discounts & gentle nudges\n\n"
    "Pain point: No real-time feedback on how much they've taken\n"
    "Value: AI-powered awareness + economic incentive (discounts)")

add_card(slide, 7.2, 1.6, 5.8, 3.5,
    "👨‍🍳 Secondary User: Hotel F&B Staff",
    "• Buffet manager / kitchen supervisor\n"
    "• Needs to know which leftover dishes are available\n"
    "• Wants to sell surplus food at discount\n"
    "• Under pressure to meet sustainability targets\n\n"
    "Pain point: Manual inventory of leftover food is slow & inaccurate\n"
    "Value: AI identification → smart meal combos → secondary sales")

add_body_text(slide, "Scenario: A typical dinner buffet at a Macau integrated resort | 200+ diners | 6 buffet stations",
              top=5.5, size=14, color=MED_GRAY)
add_page_number(slide, 3)


# ================================================================
# SLIDE 4: Solution Overview
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Our Solution: 3 AI-Powered Functions", "Prevention + Management + Reuse = Zero Waste")

add_card(slide, 0.8, 1.6, 3.8, 4.2,
    "🔥 Function 1\nCalorie Tracker",
    "📸 Take photo of food on plate\n"
    "🤖 AI identifies dishes + calories\n"
    "📊 Accumulates across multiple trips\n"
    "🐉 Cute reminder at 1800 kcal threshold\n\n"
    "Value: Prevents over-taking\n"
    "before it happens",
    title_color=ACCENT)

add_card(slide, 5.0, 1.6, 3.8, 4.2,
    "🍽️ Function 2\nWaste Scorer",
    "📸 Photograph plate after eating\n"
    "🤖 AI detects leftover food amount\n"
    "📊 Score: 0-100 (higher = less waste)\n"
    "💰 85% discount for 'Clean Plate Hero'\n\n"
    "Value: Economic incentive\n"
    "to reduce plate waste",
    title_color=GREEN)

add_card(slide, 9.2, 1.6, 3.8, 4.2,
    "🥘 Function 3\nLeftover Matcher",
    "📸 Staff photographs buffet stations\n"
    "🤖 AI identifies all remaining dishes\n"
    "👤 Diner selects dietary preference\n"
    "🎯 AI recommends 3 meal combos\n\n"
    "Value: Leftovers get\n"
    "a second life at discount",
    title_color=BLUE)

add_body_text(slide, "All three functions share the same AI engine — a closed-set matching system trained on 12 reference dishes",
              top=6.2, size=14, color=MED_GRAY)
add_page_number(slide, 4)


# ================================================================
# SLIDE 5: Value Proposition
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Value Proposition", "Measurable impact for diners, hotels, and the environment")

add_card(slide, 0.8, 1.6, 3.8, 3.0,
    "💰 For the Diner",
    "• Up to 15% discount on meals\n"
    "• Real-time calorie awareness\n"
    "• Gentle, friendly reminders\n"
    "• Personalized meal recommendations\n"
    "• Better dining experience overall",
    title_color=GREEN)

add_card(slide, 5.0, 1.6, 3.8, 3.0,
    "🏨 For the Hotel",
    "• 30% reduction in food waste\n"
    "• Lower food cost & disposal fees\n"
    "• Revenue from leftover meal combos\n"
    "• Sustainability KPI improvement\n"
    "• Enhanced brand image for ESG",
    title_color=BLUE)

add_card(slide, 9.2, 1.6, 3.8, 3.0,
    "🌍 For the Planet",
    "• Less food in landfills\n"
    "• Reduced methane emissions\n"
    "• Lower carbon footprint per meal\n"
    "• Aligned with UN SDG 12.3\n"
    "  (halve food waste by 2030)",
    title_color=RGBColor(0x8E, 0x44, 0xAD))

# Stats row
stats = [("-30%", "Food Waste"), ("85%", "Max Discount"), ("12", "Dish Library"), ("1800", "Cal Threshold")]
for i, (num, label) in enumerate(stats):
    x = 2.2 + i * 2.5
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.0), Inches(2.0), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(36)
    p.font.color.rgb = ACCENT2
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(13)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

add_page_number(slide, 5)


# ================================================================
# SLIDE 6: AI Architecture Overview (Zhang Zhenghao)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "How It Works: AI Architecture", "Multi-layer engine: DeepSeek → Perceptual Hash → PIL Fallback")

# Flowchart using shapes
layers = [
    ("📸 User Photo", "Mobile upload\nvia web app", ACCENT),
    ("🔀 AI Router", "Auto-selects\nbest engine", ACCENT2),
    ("🤖 DeepSeek Vision", "Primary: API call\nHigh accuracy", GREEN),
    ("🔍 Perceptual Hash", "Fallback: 8×8 fingerprint\nZero cost", BLUE),
    ("🎨 PIL Analysis", "Last resort: Color\nHonest but basic", MED_GRAY),
]
for i, (title, desc, color) in enumerate(layers):
    x = 0.5 + i * 2.5
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.2), Inches(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK
    p2.alignment = PP_ALIGN.CENTER
    # Arrow between boxes
    if i < len(layers) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.25), Inches(2.8), Inches(0.22), Inches(0.22))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

# Output
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "→ Output: { dish_name, calories, confidence_level }  — Always honest, always evidence-based"
p.font.size = Pt(18)
p.font.color.rgb = ACCENT
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

# Benefits
add_card(slide, 0.8, 5.8, 5.8, 1.2,
    "✅ Advantages of Multi-Layer Design",
    "• Works with or without API key  • No internet? Still works (hash + PIL)  • Never makes up fake dishes  • Graceful degradation",
    title_color=GREEN)
add_card(slide, 7.2, 5.8, 5.8, 1.2,
    "🛠️ Tech Stack",
    "Python Flask • DeepSeek Chat API • Anthropic Claude (backup) • PIL/Pillow • Perceptual Hash (custom) • Bootstrap 5",
    title_color=BLUE)

add_page_number(slide, 6)


# ================================================================
# SLIDE 7: AI Engine Details
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "AI Engine Deep Dive", "Three layers, each with a specific job")

engines = [
    ("🤖 DeepSeek Vision API", ACCENT,
     "• Sends photo as base64 + structured prompt\n"
     "• Returns JSON: dish name, category, calories, confidence\n"
     "• Model: deepseek-chat (OpenAI-compatible)\n"
     "• Config: temperature=0.1, max_tokens=2000\n"
     "• Pros: Very accurate, understands Chinese dishes\n"
     "• Cons: Costs ~$0.001 per API call"),
    ("🔍 Perceptual Hash Matcher", BLUE,
     "• Custom algorithm — zero external dependencies\n"
     "• 8×8 grayscale fingerprint + Hamming distance\n"
     "• Also: color histogram (30%) + edge texture (10%)\n"
     "• Weighted scoring: 60% hash + 30% color + 10% edge\n"
     "• Threshold: 55% minimum similarity\n"
     "• Pros: Free, offline, always available\n"
     "• Cons: Less accurate than AI, needs reference images"),
    ("🎨 PIL Color Analysis (Fallback)", MED_GRAY,
     "• Traditional image processing — no AI, no API\n"
     "• HSV color space: red-brown→meat, green→vegetables\n"
     "• 10 color-to-food mapping rules\n"
     "• Detects plate coverage ratio\n"
     "• Pros: Always works, zero cost, honest\n"
     "• Cons: Low precision, can't identify specific dishes"),
]
for i, (title, color, body) in enumerate(engines):
    y = 1.6 + i * 1.9
    add_card(slide, 0.8, y, 11.8, 1.75, title, body, title_color=color)

add_page_number(slide, 7)


# ================================================================
# SLIDE 8: Two Core Functions — now with placeholder for screenshots
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Function 1+2: Calorie Tracker & Waste Scorer", "Before-meal prevention + After-meal management")

add_card(slide, 0.8, 1.6, 5.8, 2.0,
    "🔥 Calorie Tracker (Before Meal)",
    "• User takes photo of food ON the plate\n"
    "• AI identifies each dish → estimates calories\n"
    "• Calories accumulate across multiple buffet trips\n"
    "• Progress ring shows: X / 1800 kcal\n"
    "• At threshold: cute character pops up with friendly reminder\n"
    "• Behavioral design: gentle nudge, not strict dieting",
    title_color=ACCENT)

add_card(slide, 7.2, 1.6, 5.8, 2.0,
    "🍽️ Waste Scorer (After Meal)",
    "• User takes photo of plate AFTER eating\n"
    "• AI detects leftover food amount & type\n"
    "• Weighted scoring (meat waste = 1.5× worse than vegetables)\n"
    "• Score 95+ = 'Clean Plate Hero' = 15% discount\n"
    "• Score <30 = warning message\n"
    "• Economic incentive to finish your food",
    title_color=GREEN)

# Screenshot placeholders
for i, (label, color) in enumerate([("Calorie Tracker UI", ACCENT), ("Waste Scorer UI", GREEN)]):
    x = 0.8 + i * 6.4
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.9), Inches(5.8), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    shape.line.dash_style = 2  # dashed
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"📸 [Screenshot: {label}]\n\nOpen https://buffet-zero-waste-ai.onrender.com/product1"
    p.font.size = Pt(14)
    p.font.color.rgb = MED_GRAY
    p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 8)


# ================================================================
# SLIDE 9: Prompt Evidence
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "AI Evidence: Prompt Engineering", "Actual DeepSeek API prompt, parameters, and input-output example")

# Left: Prompt code
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.8), Inches(5.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)

prompt_lines = [
    ("# System Prompt (DeepSeek Vision)", ACCENT2),
    ("", WHITE),
    ('messages = [{', WHITE),
    ('    "role": "user",', WHITE),
    ('    "content": [', WHITE),
    ('        {"type": "text", "text":', WHITE),
    ('            "You are a professional food', WHITE),
    ('             recognition assistant. This photo', WHITE),
    ('             shows a diner\'s plate at a buffet.', WHITE),
    ('             Identify remaining food items and', WHITE),
    ('             estimate amounts. Return pure JSON."', WHITE),
    ('        },', WHITE),
    ('        {"type": "image_url",', WHITE),
    ('         "image_url": {"url": "data:image/...}}', WHITE),
    ('    ]', WHITE),
    ('}]', WHITE),
    ('', WHITE),
    ("# Output Requirement:", GREEN),
    ('{"plate_status": "empty|light|moderate|...",', WHITE),
    (' "overall_waste_percentage": number,', WHITE),
    (' "items": [{"name": "dish name",', WHITE),
    ('   "category": "meat|seafood|vegetable|...",', WHITE),
    ('   "estimated_remaining_percentage": number}]}', WHITE),
    ('', WHITE),
    ('# Key: "Only list foods you ACTUALLY SEE. Be honest."', ACCENT),
]
for i, (text, color) in enumerate(prompt_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.name = 'Consolas'

# Right: API config + Example
add_card(slide, 7.8, 1.5, 5.0, 2.0,
    "⚙️ API Configuration",
    "Model: deepseek-chat\n"
    "Base URL: api.deepseek.com\n"
    "Interface: OpenAI-compatible SDK\n"
    "Temperature: 0.1 (consistent, not creative)\n"
    "Max Tokens: 2,000",
    title_color=BLUE)

add_card(slide, 7.8, 3.8, 5.0, 3.0,
    "📸 Input → Output Example",
    "INPUT: Photo of Guo Bao Rou (锅包肉)\n"
    "  → Golden crispy pork slices on white plate\n\n"
    "OUTPUT (JSON):\n"
    '{\n'
    '  "plate_status": "full",\n'
    '  "items": [{\n'
    '    "name": "锅包肉",\n'
    '    "category": "肉类",\n'
    '    "estimated_remaining_percentage": 85\n'
    '  }],\n'
    '  "summary": "盘中主要是锅包肉，约剩余85%"\n'
    '}',
    title_color=GREEN)

add_page_number(slide, 9)


# ================================================================
# SLIDE 10: Code Evidence
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "AI Evidence: Code & Implementation", "Key code snippets from the actual project")

# Left: DeepSeekVision code
code_shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(6.0), Inches(2.8))
code_shape1.fill.solid()
code_shape1.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
code_shape1.line.fill.background()
tf1 = code_shape1.text_frame
tf1.word_wrap = True
tf1.margin_left = Inches(0.25)
tf1.margin_top = Inches(0.15)

lines1 = [
    ("# DeepSeekVision Class (app.py line 793)", ACCENT2),
    ("class DeepSeekVision:", WHITE),
    ('    def __init__(self):', WHITE),
    ('        self.api_key = os.environ.get("DEEPSEEK_API_KEY")', WHITE),
    ('        from openai import OpenAI', WHITE),
    ('        self.client = OpenAI(', WHITE),
    ('            api_key=self.api_key,', WHITE),
    ('            base_url="https://api.deepseek.com"  ', GREEN),
    ('        )', WHITE),
    ('', WHITE),
    ('    def analyze_plate_waste(self, image_path):', WHITE),
    ('        img_b64 = base64.b64encode(...)', WHITE),
    ('        resp = self.client.chat.completions.create(', WHITE),
    ('            model="deepseek-chat",', WHITE),
    ('            messages=[{...image + prompt...}],', WHITE),
    ('            max_tokens=2000, temperature=0.1', WHITE),
    ('        )', WHITE),
]
for i, (text, color) in enumerate(lines1):
    if i == 0:
        p = tf1.paragraphs[0]
    else:
        p = tf1.add_paragraph()
    p.text = text
    p.font.size = Pt(10.5)
    p.font.color.rgb = color
    p.font.name = 'Consolas'

# Right: Perceptual Hash code
code_shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(6.0), Inches(2.8))
code_shape2.fill.solid()
code_shape2.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
code_shape2.line.fill.background()
tf2 = code_shape2.text_frame
tf2.word_wrap = True
tf2.margin_left = Inches(0.25)
tf2.margin_top = Inches(0.15)

lines2 = [
    ("# Perceptual Hash Algorithm (app.py line 1132)", ACCENT2),
    ("def _compute_ahash(img):", WHITE),
    ("    gray = img.convert('L').resize((8, 8))", WHITE),
    ("    pixels = list(gray.getdata())", WHITE),
    ("    avg = sum(pixels) / 64", WHITE),
    ("    return ''.join(              ", WHITE),
    ("        '1' if p > avg else '0'  ", GREEN),
    ("        for p in pixels)", WHITE),
    ("", WHITE),
    ("def _hamming_distance(h1, h2):", WHITE),
    ("    return sum(c1 != c2             ", WHITE),
    ("        for c1, c2 in zip(h1, h2))", WHITE),
    ("", WHITE),
    ("# Similarity = 1.0 - distance/64", WHITE),
    ("# Higher = more likely same dish", WHITE),
]
for i, (text, color) in enumerate(lines2):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = text
    p.font.size = Pt(10.5)
    p.font.color.rgb = color
    p.font.name = 'Consolas'

# Bottom: Input-output examples
add_card(slide, 0.5, 4.6, 3.8, 2.5,
    "📸 Input: Photo",
    "User uploads a plate photo\n"
    "→ Guo Bao Rou (锅包肉)\n"
    "→ Golden crispy pork slices\n"
    "→ On a white round plate\n\n"
    "[Photo sent as base64 to API]",
    title_color=ACCENT)

add_card(slide, 4.6, 4.6, 3.8, 2.5,
    "🤖 AI Processing",
    "Engine tries in order:\n"
    "1. DeepSeek Vision API\n"
    "2. Perceptual Hash Match\n"
    "3. PIL Color Analysis\n\n"
    "Each returns structured JSON\n"
    "with confidence score",
    title_color=BLUE)

add_card(slide, 8.7, 4.6, 4.3, 2.5,
    "📊 Output: Result",
    "Dish Name: 锅包肉\n"
    "Category: Meat (肉类)\n"
    "Estimated Calories: 310 kcal\n"
    "Confidence: High (78%)\n"
    "Evidence: Golden-orange irregular\n"
    "pieces, glossy sauce coating",
    title_color=GREEN)

add_page_number(slide, 10)


# ================================================================
# SLIDE 11: Iteration Log
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Iteration Log: v1.0 → v2.5", "6 versions, 6 days — what we fixed and why")

# Timeline
iterations = [
    ("v1.0", "Jun 10\nNight", "Random AI generation\n❌ Fake results", "User said 'This is fraud'\n→ Must change"),
    ("v2.0", "Jun 11\nMorning", "Fixed dish library\n✅ Closed-set matching", "AI only matches 12\nknown dishes"),
    ("v2.1→v2.3", "Jun 11-15\nDaytime", "Calorie tracking\n✅ Perceptual hash\n✅ Base64 storage", "Multi-feature fusion\nImages survive deploys"),
    ("v2.4", "Jun 15\nEvening", "DeepSeek API\n✅ Demo mode", "Much higher accuracy\nClassroom-ready demo"),
    ("v2.5", "Jun 15\nNight", "P1-P4 improvements\n✅ Allergy filter\n✅ Real scoring", "Vegans never get meat\n3 truly different combos"),
]
for i, (version, date, features, fix) in enumerate(iterations):
    x = 0.3 + i * 2.6
    # Version header
    add_card(slide, x, 1.5, 2.4, 1.4, f"{version}\n{date}", "", title_color=ACCENT2)
    # Features
    add_card(slide, x, 3.1, 2.4, 2.2, "New Features", features, title_color=GREEN)
    # Fixes
    add_card(slide, x, 5.5, 2.4, 1.6, "Key Fix", fix, title_color=ACCENT)
    # Arrow
    if i < len(iterations) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.45), Inches(2.5), Inches(0.12), Inches(0.12))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MED_GRAY
        arrow.line.fill.background()

add_page_number(slide, 11)


# ================================================================
# SLIDE 12: AI Evidence Summary (Team Member D)
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Evidence Summary: Proof of AI Work", "All the receipts — prompts, code, API config, iteration history")

items = [
    ("✅ Prompt Engineering", "Custom prompts with role-setting, output formatting, and honesty constraints"),
    ("✅ API Configuration", "DeepSeek Chat API: temperature=0.1, max_tokens=2000, OpenAI-compatible SDK"),
    ("✅ Code Implementation", "DeepSeekVision class + custom perceptual hash + multi-feature SmartMatcher"),
    ("✅ Input-Output Pairs", "Photo → JSON: dish name, category, calories, confidence — documented examples"),
    ("✅ Iteration History", "v1.0 (random) → v2.5 (DeepSeek+P1-P4): 6 versions in 6 days"),
    ("✅ Live Demo Ready", "Fully functional web app at buffet-zero-waste-ai.onrender.com"),
]
for i, (label, desc) in enumerate(items):
    y = 1.5 + i * 0.95
    add_card(slide, 0.8, y, 11.8, 0.8, label, desc, title_color=GREEN)

# GitHub link
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5 + 6 * 0.95 + 0.1), Inches(11.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "🔗 GitHub: https://github.com/986867562zzh-hub/buffet-zero-waste-ai  |  Live: https://buffet-zero-waste-ai.onrender.com"
p.font.size = Pt(14)
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

add_page_number(slide, 12)


# ================================================================
# SLIDE 13-15: Live Demo placeholder slides
# ================================================================
for demo_num, demo_title, demo_desc in [
    (13, "Live Demo: Calorie Tracking (Steps 1-4)",
     "🔥 Taking photos of 4 dishes → AI identifies each one → Calories accumulate → Progress ring fills up"),
    (14, "Live Demo: Threshold Alert (Steps 5-8)",
     "🐉 Calories reach 1800+ kcal → Cute cartoon character pops up → Gentle reminder to stop taking more"),
    (15, "Live Demo: Leftover Meal Matching",
     "🥘 Staff uploads buffet photos → AI identifies dishes → Select dietary type → Get 3 personalized combos"),
]:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, demo_title, demo_desc)

    # Big demo area
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.0), Inches(10.3), Inches(4.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    shape.line.color.rgb = ACCENT2
    shape.line.width = Pt(3)
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"🎬 LIVE DEMO — Step {demo_num-12}"
    p.font.size = Pt(28)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "\n\nOpen: https://buffet-zero-waste-ai.onrender.com/demo"
    p2.font.size = Pt(20)
    p2.font.color.rgb = BLUE
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "\nPresenter: [Team Member C]"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MED_GRAY
    p3.alignment = PP_ALIGN.CENTER

    add_page_number(slide, demo_num)


# ================================================================
# SLIDE 16: Summary & Business Value
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Summary: What We Built & Why It Matters", "A practical AI solution for a real hospitality problem")

add_card(slide, 0.8, 1.5, 5.8, 2.5,
    "🎯 The Problem We Solved",
    "• Buffet food waste: a measurable, daily problem\n"
    "• Two pain points: over-taking (consumer) + inventory (hotel)\n"
    "• Existing solutions: manual, slow, or non-existent\n\n"
    "Our approach: Make AI do the recognition work,\n"
    "then use behavioral design + economic incentives",
    title_color=ACCENT)

add_card(slide, 7.2, 1.5, 5.8, 2.5,
    "💡 What Makes It Smart",
    "• Closed-set matching: AI only works with known dishes\n"
    "  → No hallucination, no fake results\n"
    "• Multi-layer engine: works with or without API key\n"
    "• Gentle behavioral intervention (not strict rules)\n"
    "• Real, working prototype — not just slides",
    title_color=BLUE)

add_card(slide, 0.8, 4.3, 5.8, 2.0,
    "📊 Measurable Impact (Projected)",
    "• 30% reduction in consumer plate waste\n"
    "• 20% additional revenue from leftover meal combos\n"
    "• 15% increase in customer satisfaction (discounts + personalization)\n"
    "• Aligned with Macau's green hotel certification standards",
    title_color=GREEN)

add_card(slide, 7.2, 4.3, 5.8, 2.0,
    "🛠️ Tech That Actually Works",
    "• DeepSeek Vision API — cheap, accurate, Chinese-fluent\n"
    "• Perceptual hash — custom algorithm, zero API cost\n"
    "• 12-dish reference library — base64 stored, never lost\n"
    "• 6 versions in 6 days — iterative, evidence-based development",
    title_color=ACCENT2)

add_page_number(slide, 16)


# ================================================================
# SLIDE 17: Future Roadmap + Thank You
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "Future Roadmap & Thank You", "Where this project goes next")

add_card(slide, 0.8, 1.5, 3.8, 2.5,
    "🚀 Phase 2 (Next Month)",
    "• Expand dish library: 12 → 50 dishes\n"
    "• Multi-angle reference photos (3-5 per dish)\n"
    "• Train a custom image classification model\n"
    "• Mobile app (React Native)\n"
    "• QR code integration for easy access",
    title_color=BLUE)

add_card(slide, 5.0, 1.5, 3.8, 2.5,
    "🏨 Phase 3 (Semester 2)",
    "• Hotel PMS integration (Opera, OnQ)\n"
    "• Real-time kitchen display for staff\n"
    "• Dynamic pricing based on shelf life\n"
    "• Customer preference learning\n"
    "• A/B testing with partner hotel",
    title_color=ACCENT2)

add_card(slide, 9.2, 1.5, 3.8, 2.5,
    "🌍 Long-term Vision",
    "• Deploy in Macau integrated resorts\n"
    "• Multi-language support (CN/EN/PT/KR)\n"
    "• Carbon footprint tracking per dish\n"
    "• Open-source the dish library format\n"
    "• Publish findings in hospitality journal",
    title_color=GREEN)

# Thank you section
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.8), Inches(2.2))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
p.text = "Thank You! 🙏"
p.font.size = Pt(36)
p.font.color.rgb = ACCENT2
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "\nQuestions & Discussion"
p2.font.size = Pt(22)
p2.font.color.rgb = WHITE
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "🔗 GitHub: github.com/986867562zzh-hub/buffet-zero-waste-ai  |  🌐 Live: buffet-zero-waste-ai.onrender.com"
p3.font.size = Pt(13)
p3.font.color.rgb = MED_GRAY
p3.alignment = PP_ALIGN.CENTER

add_page_number(slide, 17)


# ===== SAVE =====
output_path = r'D:\Quant_Projects\Quant_Projects\buffet-demo\Group_Project_Presentation.pptx'
prs.save(output_path)
print(f'✅ PPT saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
